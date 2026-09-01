#!/usr/bin/env python3
"""既存デッキの体裁を実測する。手を動かす前に必ず通す。

    python3 pptx_measure.py deck.pptx --slides 4,6,9
    python3 pptx_measure.py deck.pptx            # 色と線幅の集計のみ

依存: python-pptx, lxml
"""
import argparse
import collections
import zipfile

from lxml import etree
from pptx import Presentation

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
IN = 914400


def measure_slides(path, indices):
    p = Presentation(path)
    for idx in indices:
        if idx < 1 or idx > len(p.slides):
            print("スライド%d は存在しません" % idx)
            continue
        s = p.slides[idx - 1]
        print("=" * 60, "SLIDE", idx)
        for sh in s.shapes:
            if sh.left is None:
                continue
            pg = sh._element.find(".//" + A + "prstGeom")
            info = {}
            if sh.has_text_frame:
                bp = sh.text_frame._txBody.find(A + "bodyPr")
                info["sz"] = sorted({r.font.size.pt for pa in sh.text_frame.paragraphs
                                     for r in pa.runs if r.font.size})
                info["color"] = sorted({str(r.font.color.rgb) for pa in sh.text_frame.paragraphs
                                        for r in pa.runs
                                        if r.font.color and r.font.color.type == 1})
                info["font"] = sorted({r.font.name for pa in sh.text_frame.paragraphs
                                       for r in pa.runs if r.font.name})
                if bp is not None:
                    info["anchor"] = bp.get("anchor")
                    info["tIns"] = bp.get("tIns")
            print("%-28r geo=%-10s (%.2f,%.2f,%.2f,%.2f) %s" % (
                sh.name, pg.get("prst") if pg is not None else "-",
                sh.left / IN, sh.top / IN, sh.width / IN, sh.height / IN, info))
            if sh.has_text_frame and sh.text_frame.text.strip():
                print("     TXT:", sh.text_frame.text.replace("\n", " | ")[:80])


def measure_palette(path):
    col, wid = collections.Counter(), collections.Counter()
    z = zipfile.ZipFile(path)
    for n in z.namelist():
        if not n.startswith("ppt/slides/slide") or not n.endswith(".xml"):
            continue
        x = etree.fromstring(z.read(n))
        for e in x.iter(A + "srgbClr"):
            col[e.get("val")] += 1
        for ln in x.iter(A + "ln"):
            sf = ln.find(A + "solidFill/" + A + "srgbClr")
            if ln.get("w") and sf is not None:
                wid[(ln.get("w"), sf.get("val"))] += 1
    z.close()
    print("\n色（出現頻度順）:", col.most_common(20))
    print("線幅(EMU, 12700=1pt):", wid.most_common(10))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--slides", default="", help="実測するスライド番号（1始まり、カンマ区切り）")
    args = ap.parse_args()
    if args.slides:
        measure_slides(args.deck, [int(x) for x in args.slides.split(",") if x.strip()])
    measure_palette(args.deck)


if __name__ == "__main__":
    main()
