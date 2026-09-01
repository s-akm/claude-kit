#!/usr/bin/env python3
"""提出前の PowerPoint を機械で検査する。元ファイルは変更しない。

    python3 pptx_qa.py 説明資料.pptx --out-dir 90_qa/2026-09-01-1200

検査する内容
    1. 文字切れ（推定した本文の高さ > 枠の高さ）
    2. フォントの混在
    3. 図形の重なり
    4. スライド番号の欠落・重複
    5. タイトルの重複
    6. 出典の記載（数値や図表があるスライドに出典が無い）
    7. 空のテキスト枠

before/after の文言差分や書式の一様性は pptx_audit.py が担当する。
この検査は「提出前に 1 ファイルだけ見る」用途に絞っている。

終了コード
    0 指摘なし / 1 警告あり / 2 提出できない指摘あり、またはファイルが読めない
"""
import argparse
import os
import re
import sys
from collections import Counter, defaultdict

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import qa_report  # noqa: E402
from pptx_audit import est_height_in, walk  # noqa: E402

from pptx import Presentation  # noqa: E402

IN = 914400
NUM_ONLY = re.compile(r"^\s*\d{1,3}\s*$")
FIGURE_WORDS = re.compile(r"図\s?\d|表\s?\d|グラフ|チャート|[\d.]+\s?[%％]|億円|万円")
SOURCE_WORDS = re.compile(r"出典|出所|Source|ソース|注\)|※")


def add(f, sev, check, where, what, how=""):
    f.append({"重要度": sev, "検査": check, "場所": where, "内容": what, "対応": how})


def rect(sh):
    try:
        return (sh.left / IN, sh.top / IN, sh.width / IN, sh.height / IN)
    except TypeError:
        return None


def overlap_ratio(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    inter = ix * iy
    if inter <= 0:
        return 0.0
    smaller = min(aw * ah, bw * bh)
    return inter / smaller if smaller > 0 else 0.0


def check(path, overlap_threshold=0.15):
    prs = Presentation(path)
    findings = []
    notes = []
    titles = defaultdict(list)
    page_numbers = []
    fonts = Counter()

    for idx, slide in enumerate(prs.slides, start=1):
        where = "スライド %d" % idx
        boxes = []
        top_texts = []

        for sh in walk(slide.shapes):
            if not getattr(sh, "has_text_frame", False):
                continue
            text = sh.text_frame.text.strip()
            r = rect(sh)

            for p in sh.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.name:
                        fonts[run.font.name] += 1

            if not text:
                if r and r[2] > 0.3 and r[3] > 0.2:
                    add(findings, "情報", "空のテキスト枠", where,
                        "文字の入っていない枠がある（左 %.1fin 上 %.1fin）" % (r[0], r[1]),
                        "使っていない枠なら消す")
                continue

            if NUM_ONLY.match(text):
                page_numbers.append((idx, text.strip()))
                continue

            if r:
                boxes.append((sh, r, text))
                top_texts.append((r[1], text))

                # --- 文字切れ
                try:
                    need = est_height_in(sh)
                except Exception:
                    need = None
                if need and r[3] > 0 and need > r[3] * 1.02:
                    add(findings, "警告", "文字切れ", where,
                        "枠の高さ %.2fin に対して本文が %.2fin 必要「%s」"
                        % (r[3], need, text[:28]),
                        "枠を広げるか、文字を減らす")

        # --- 図形の重なり
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                ratio = overlap_ratio(boxes[i][1], boxes[j][1])
                if ratio >= overlap_threshold:
                    add(findings, "警告", "図形の重なり", where,
                        "「%s」と「%s」が %.0f%% 重なっている"
                        % (boxes[i][2][:18], boxes[j][2][:18], ratio * 100),
                        "位置をずらすか、片方を消す")

        # --- タイトル（そのスライドで一番上にある文字）
        if top_texts:
            top_texts.sort()
            titles[top_texts[0][1]].append(idx)

        # --- 出典
        body = " ".join(t for _, t in top_texts)
        if FIGURE_WORDS.search(body) and not SOURCE_WORDS.search(body):
            add(findings, "情報", "出典", where,
                "数値か図表があるが、出典の記載が見あたらない",
                "出典が要るスライドか目視で判断する")

    # --- タイトル重複
    for t, pages in titles.items():
        if len(pages) > 1:
            add(findings, "警告", "タイトル重複", "スライド " + ",".join(map(str, pages)),
                "同じ見出し「%s」が %d 枚にある" % (t[:30], len(pages)),
                "続きのスライドなら「（続き）」を付けるなどで区別する")

    # --- スライド番号
    total = len(prs.slides)
    numbered = {i for i, _ in page_numbers}
    missing = [i for i in range(1, total + 1) if i not in numbered]
    if missing and len(missing) == total:
        add(findings, "警告", "スライド番号", "全体",
            "スライド番号が 1 枚も見あたらない（%d 枚）" % total,
            "マスターにページ番号を入れるか、意図して外しているか確認する")
    elif missing:
        add(findings, "警告", "スライド番号", "スライド " + ",".join(map(str, missing[:20])),
            "番号が見あたらないスライドが %d 枚ある" % len(missing),
            "抜けているページに番号を入れる")
    dup = [n for n, c in Counter(n for _, n in page_numbers).items() if c > 1]
    if dup:
        add(findings, "警告", "スライド番号", "全体",
            "同じ番号が複数のスライドにある: %s" % ",".join(dup[:10]),
            "手で書いた番号が残っていないか確認する")

    # --- フォント混在
    if len(fonts) > 1:
        add(findings, "警告", "フォント混在", "全体",
            "本文に %d 種類のフォントが使われている（%s）"
            % (len(fonts), " / ".join("%s×%d" % kv for kv in fonts.most_common(5))),
            "テーマとマスターを含めた全レイヤーでそろえる。fix_fonts.py を使う")

    notes.append("スライド番号は「数字だけの文字枠」で判定している。"
                 "マスター側のプレースホルダで表示している場合は検出できない。")
    notes.append("文字切れは推定であり、実際の描画とは差が出る。"
                 "確定させるには PowerPoint から PDF に書き出して目視する。")
    notes.append("出典の要否と、Excel との数値の一致は、この検査では判定しない。")
    return findings, notes


def main():
    ap = argparse.ArgumentParser(description="PowerPoint の機械検査（元ファイルは変更しない）")
    ap.add_argument("file", help="検査する .pptx")
    ap.add_argument("--out-dir", default=".", help="結果の出力先")
    ap.add_argument("--stem", help="出力ファイル名（既定: pptx-<ファイル名>）")
    ap.add_argument("--overlap", type=float, default=0.15,
                    help="重なりとみなす面積比（既定 0.15）")
    args = ap.parse_args()

    if not os.path.isfile(args.file):
        print("ファイルがありません: %s" % args.file, file=sys.stderr)
        return 2

    versions = qa_report.tool_versions(["pptx"])
    try:
        findings, notes = check(args.file, args.overlap)
    except Exception as e:
        print("読めませんでした: %s: %s" % (type(e).__name__, e), file=sys.stderr)
        return 2

    report = qa_report.build("PowerPoint", args.file, findings, versions, notes)
    stem = args.stem or ("pptx-" + os.path.splitext(os.path.basename(args.file))[0])
    jp, mp = qa_report.write(report, args.out_dir, stem)
    c = report["件数"]
    print("PPTX   %s  要確認 %d / 警告 %d / 情報 %d  → %s"
          % (os.path.basename(args.file), c["要確認"], c["警告"], c["情報"], mp))
    return qa_report.exit_code(report)


if __name__ == "__main__":
    sys.exit(main())
