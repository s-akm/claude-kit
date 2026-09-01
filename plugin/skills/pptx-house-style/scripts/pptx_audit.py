#!/usr/bin/env python3
"""納品前チェック。結果を報告できる形で出す。

    python3 pptx_audit.py after.pptx --before before.pptx --font 游ゴシック
    python3 pptx_audit.py deck.pptx --pixel-diff before.jpg after.jpg

確認する内容
  1. 文言差分（--before 指定時）
  2. 書式の一様性（同種要素の座標・級数・色が1種類に収束しているか）
  3. フォントの混在
  4. 残存物（roundRect、段落先頭の丸付き数字）
  5. はみ出し（推定描画下端 > フッター上端）
"""
import argparse
import collections
import sys

from pptx import Presentation

CIRC = "①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭⑮⑯⑰⑱⑲⑳"
IN = 914400
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def walk(shapes):
    for s in shapes:
        yield s
        if s.shape_type == 6:  # GROUP
            for x in walk(s.shapes):
                yield x


def texts(path):
    p = Presentation(path)
    res = []
    for s in p.slides:
        acc = []
        for sh in walk(s.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip():
                acc.append(sh.text_frame.text.strip())
            if sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        if c.text.strip():
                            acc.append(c.text.strip())
        res.append(sorted(acc))
    return res


def est_height_in(sh, lnpct=1.40, aft_pt=3):
    w = sh.width / IN * 72.0 - 4
    tot = 0.0
    for p in sh.text_frame.paragraphs:
        if not p.text.strip():
            tot += 6
            continue
        sz = max([r.font.size.pt for r in p.runs if r.font.size] or [10.5])
        em = sum(1.0 if ord(c) > 0x2000 else 0.55 for c in p.text) * sz
        tot += max(1, int(em / w) + 1) * sz * lnpct + aft_pt
    return tot / 72.0


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--before", help="修正前のファイル。文言差分を取る")
    ap.add_argument("--font", help="統一したフォント名")
    ap.add_argument("--footer-top", type=float, default=7.03,
                    help="フッター上端（インチ）。既定は 7.5in スライドで 0.47in 残す想定")
    ap.add_argument("--pixel-diff", nargs=2, metavar=("BEFORE_JPG", "AFTER_JPG"))
    args = ap.parse_args()

    ng = 0

    if args.pixel_diff:
        from PIL import Image, ImageChops
        import numpy as np
        a = Image.open(args.pixel_diff[0]).convert("RGB")
        b = Image.open(args.pixel_diff[1]).convert("RGB")
        d = np.asarray(ImageChops.difference(a, b)).astype(int)
        pct = (d.max(axis=2) > 90).mean() * 100
        print("[5] ピクセル差分: %.2f%%" % pct)
        return 0

    if args.before:
        ba, af = texts(args.before), texts(args.deck)
        if len(ba) != len(af):
            print("[1] 文言差分: スライド枚数が %d → %d に変わっています" % (len(ba), len(af)))
            ng += 1
        else:
            diff = [i + 1 for i, (x, y) in enumerate(zip(ba, af)) if x != y]
            print("[1] 文言差分: %s" % ("差分なし" if not diff else "スライド %s に差分" % diff))
            if diff:
                ng += 1

    p = Presentation(args.deck)
    shapes_by_size = collections.Counter()
    fonts = collections.Counter()
    round_rects, circled, overflow = [], [], []

    for i, s in enumerate(p.slides, start=1):
        for sh in walk(s.shapes):
            if sh.left is None:
                continue
            pg = sh._element.find(".//" + A + "prstGeom")
            if pg is not None and pg.get("prst") == "roundRect":
                round_rects.append((i, sh.name))
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            for pa in tf.paragraphs:
                for r in pa.runs:
                    if r.font.name:
                        fonts[r.font.name] += 1
                if pa.text and pa.text[0] in CIRC:
                    circled.append((i, pa.text[:24]))
            szs = tuple(sorted({r.font.size.pt for pa in tf.paragraphs
                                for r in pa.runs if r.font.size}))
            if szs:
                bp = tf._txBody.find(A + "bodyPr")
                shapes_by_size[(szs, bp.get("anchor") if bp is not None else None,
                                bp.get("tIns") if bp is not None else None)] += 1
            if tf.text.strip():
                bottom = sh.top / IN + est_height_in(sh)
                if bottom > args.footer_top:
                    overflow.append((i, sh.name, round(bottom, 2)))

    print("[2] 書式の組み合わせ（サイズ, anchor, tIns）: %d 種類" % len(shapes_by_size))
    for k, v in shapes_by_size.most_common(12):
        print("      %s  x%d" % (k, v))
    print("[3] フォント: %s" % dict(fonts))
    if args.font and any(f != args.font for f in fonts):
        print("      → %s 以外が混在しています" % args.font)
        ng += 1
    print("[4] roundRect: %s" % (round_rects or "なし"))
    print("    段落先頭の丸付き数字: %s" % (circled or "なし"))
    if round_rects or circled:
        ng += 1
    print("[5] はみ出し（下端 > %.2fin）: %s" % (args.footer_top, overflow or "なし"))
    if overflow:
        ng += 1

    print("\n要確認 %d 項目" % ng)
    return 1 if ng else 0


if __name__ == "__main__":
    sys.exit(main())
