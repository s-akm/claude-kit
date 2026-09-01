#!/usr/bin/env python3
"""検査スクリプトの確認用サンプルを作る。

顧客のファイルは使わない。毎回ここからゼロで作る。

    uv run python tests/make_fixtures.py

作るもの（tests/fixtures/ 配下）
    ok.xlsx  / ng.xlsx   Excel の正常版と問題版
    ok.pptx  / ng.pptx   PowerPoint の正常版と問題版
    ok.pdf   / ng.pdf    PDF の正常版と問題版

ng 側に仕込む問題
    xlsx  数式エラー、外部リンク、非表示シート・行・列、壊れた名前定義、
          存在しないシートへの参照、コメントの付いていない未決セル、書式の揺れ
    pptx  枠からはみ出す本文、フォント混在、図形の重なり、タイトル重複、
          スライド番号なし
    pdf   文字が抽出できないページ
"""
import os
import shutil
import subprocess
import sys
import zipfile

from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.styles import Font
from openpyxl.workbook.defined_name import DefinedName
from openpyxl.worksheet.datavalidation import DataValidation

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "fixtures")


# ---------------------------------------------------------------- xlsx
def build_xlsx_ok(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "機能一覧"
    ws.append(["ID", "分類", "要件", "受入条件", "状態"])
    ws.append(["FR-001", "画面", "編集者は公開日時を設定できる", "未来日時の記事が一覧に出ない", "確定"])
    ws.append(["FR-002", "連携", "夜間に在庫を取り込む", "取り込み件数が記録される", "確定"])
    for c in "ABCDE":
        ws.column_dimensions[c].width = 18
    ws2 = wb.create_sheet("凡例")
    ws2.append(["記号", "意味"])
    ws2.append(["確定", "合意済み"])
    dv = DataValidation(type="list", formula1='"確定,調整中,未決"', allow_blank=True)
    ws.add_data_validation(dv)
    dv.add("E2:E100")
    wb.save(path)


def build_xlsx_ng(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "機能一覧"
    ws.append(["ID", "分類", "要件", "件数", "状態"])
    ws.append(["FR-001", "画面", "編集者は公開日時を設定できる", 100, "確定"])
    # 数式エラーになる式（0 除算）
    ws.append(["FR-002", "連携", "夜間に在庫を取り込む", "=D2/0", "要確認"])
    # 式に残った #REF!
    ws["G2"] = "=SUM(#REF!)"
    # 存在しないシートへの参照
    ws.append(["FR-003", "帳票", "月次帳票を出力する", "=存在しないシート!A1", "未決"])
    # 外部ブックへの参照
    ws.append(["FR-004", "連携", "外部の一覧を取り込む", "='C:\\tmp\\[外部ブック.xlsx]Sheet1'!A1", "確定"])
    # 書式の揺れ（同じ列に数値と文字列の日付）
    ws["F1"] = "期限"
    ws["F2"] = "2026-09-15"
    ws["F3"] = "2026/9/20"
    ws["F4"] = 46000
    ws["F4"].number_format = "yyyy/mm/dd"
    ws["F5"] = 1234.5
    ws["F5"].number_format = "#,##0.00"
    # コメントの付いていない未決セル
    ws["E3"].comment = Comment("連携先の確認待ち", "レビュー")
    # E4（未決）にはコメントを付けない

    hidden = wb.create_sheet("作業メモ")
    hidden.append(["社内メモ", "この行は提出前に消す"])
    hidden.sheet_state = "hidden"

    ws.row_dimensions[3].hidden = True
    ws.column_dimensions["C"].hidden = True

    wb.defined_names.add(DefinedName("壊れた名前", attr_text="#REF!"))
    wb.defined_names.add(DefinedName("使っていない名前", attr_text="機能一覧!$A$1"))
    wb.save(path)


# ---------------------------------------------------------------- pptx
def _txt(slide, left, top, width, height, text, size=14, font="Noto Sans CJK JP"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.name = font
    return box


def build_pptx_ok(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i, (title, body) in enumerate(
        [
            ("要件定義の進め方", "機能要件と非機能要件を分けて整理する。"),
            ("現行の課題", "公開日時を設定できないため、担当者が深夜に手作業で公開している。"),
            ("対応方針", "公開予約の仕組みを追加する。既存の権限設定はそのまま使う。"),
        ],
        start=1,
    ):
        s = prs.slides.add_slide(blank)
        _txt(s, 0.6, 0.5, 12.0, 0.8, title, size=24)
        _txt(s, 0.6, 1.6, 12.0, 3.0, body, size=14)
        _txt(s, 12.3, 6.9, 0.6, 0.35, str(i), size=10)
    prs.save(path)


def build_pptx_ng(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    _txt(s1, 0.6, 0.5, 12.0, 0.8, "対応方針", size=24)
    # 枠に対して明らかに長い本文（はみ出し）
    _txt(s1, 0.6, 1.6, 4.0, 0.6,
         "公開予約の仕組みを追加する。既存の権限設定はそのまま使う。"
         "移行は段階的に行い、既存の記事には影響を与えない。"
         "運用手順書もあわせて更新する。", size=14)
    # 重なる図形
    _txt(s1, 3.0, 1.8, 4.0, 1.0, "重なっている枠", size=14)

    s2 = prs.slides.add_slide(blank)
    # タイトル重複
    _txt(s2, 0.6, 0.5, 12.0, 0.8, "対応方針", size=24)
    # フォント混在
    _txt(s2, 0.6, 1.6, 12.0, 1.0, "別のフォントで書かれた本文", size=14, font="DejaVu Sans")
    _txt(s2, 0.6, 3.0, 12.0, 1.0, "図 1 の値は前年比で 120% になる", size=14)
    # スライド番号なし
    prs.save(path)


# ---------------------------------------------------------------- pdf
def build_pdfs(ok_pptx, ok_pdf, ng_pdf):
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        print("  soffice が無いため PDF の fixture は作りません", file=sys.stderr)
        return False
    tmp = os.path.join(OUT, "_pdftmp")
    os.makedirs(tmp, exist_ok=True)
    subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                    "--outdir", tmp, ok_pptx],
                   check=True, capture_output=True, timeout=180)
    made = os.path.join(tmp, os.path.splitext(os.path.basename(ok_pptx))[0] + ".pdf")
    shutil.copy(made, ok_pdf)

    # ng: 文字を持たないページだけの PDF（画像化したのと同じ状態を作る）
    from pypdf import PdfWriter
    w = PdfWriter()
    w.add_blank_page(width=842, height=595)
    w.add_blank_page(width=842, height=595)
    with open(ng_pdf, "wb") as f:
        w.write(f)
    shutil.rmtree(tmp, ignore_errors=True)
    return True


def main():
    os.makedirs(OUT, exist_ok=True)
    build_xlsx_ok(os.path.join(OUT, "ok.xlsx"))
    build_xlsx_ng(os.path.join(OUT, "ng.xlsx"))
    build_pptx_ok(os.path.join(OUT, "ok.pptx"))
    build_pptx_ng(os.path.join(OUT, "ng.pptx"))
    made = build_pdfs(os.path.join(OUT, "ok.pptx"),
                      os.path.join(OUT, "ok.pdf"),
                      os.path.join(OUT, "ng.pdf"))
    for n in sorted(os.listdir(OUT)):
        p = os.path.join(OUT, n)
        if os.path.isfile(p):
            print("  %-10s %7d bytes" % (n, os.path.getsize(p)))
    if not made:
        print("  PDF は未作成", file=sys.stderr)


if __name__ == "__main__":
    main()
